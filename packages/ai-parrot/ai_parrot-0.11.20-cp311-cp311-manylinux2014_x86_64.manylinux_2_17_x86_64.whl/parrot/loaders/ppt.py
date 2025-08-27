from typing import List
from pathlib import PurePath
import logging
from pptx import Presentation
from langchain.docstore.document import Document
from .abstract import AbstractLoader

class PowerPointLoader(AbstractLoader):
    """
    Load Microsoft PowerPoint .pptx files as Langchain Documents.
    Generates one Document per slide with text.
    """
    extensions: List[str] = ['.pptx']

    def extract_slide_text(self, slide):
        """Extract all text from a slide as a single string."""
        text_chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_chunks.append(shape.text.strip())
        return "\n\n".join(text_chunks).strip()

    def slide_has_text(self, slide) -> bool:
        """Determine if a slide contains any text."""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return True
        return False

    def slide_has_images_only(self, slide) -> bool:
        """Return True if slide has images and no text."""
        has_image = False
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE shape type in python-pptx
                has_image = True
            if hasattr(shape, "text") and shape.text.strip():
                return False
        return has_image

    async def _load(self, path: PurePath, **kwargs) -> List[Document]:
        """Load data from a .pptx file and return one Document per slide with text."""
        self.logger.info(f"Loading PowerPoint file: {path}")
        docs = []
        prs = Presentation(str(path))
        slide_count = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            # Ignore slides with only images and no text
            if self.slide_has_images_only(slide):
                self.logger.info(f"Slide {i+1}/{slide_count}: only images, skipping.")
                continue
            text = self.extract_slide_text(slide)
            if not text:
                # No text, no image: skip (empty slide)
                self.logger.info(f"Slide {i+1}/{slide_count}: no text content, skipping.")
                continue
            slide_meta = {
                "slide_number": i+1,
                "slide_id": slide.slide_id,
                "slide_notes": "",
            }
            # Optionally extract slide notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_meta["slide_notes"] = slide.notes_slide.notes_text_frame.text.strip()
            # Merge with global metadata
            metadata = self.create_metadata(
                path=path,
                doctype="pptx",
                source_type="powerpoint",
                doc_metadata=slide_meta,
            )
            # Slide context string (header)
            context_str = (
                f"File Name: {path.name}\n"
                f"Slide Number: {i+1}\n"
                f"Slide ID: {slide.slide_id}\n"
                f"Document Type: pptx\n"
                f"Source Type: powerpoint\n"
                "======\n"
            )
            doc = self.create_document(
                content=context_str + text,
                path=path,
                metadata=metadata
            )
            docs.append(doc)

        return docs
