"""
파일 텍스트 추출기 인터페이스
전략 패턴으로 다양한 파일 형식 지원
"""

import asyncio
from abc import ABC, abstractmethod
from typing import Optional, Dict, Any
from dataclasses import dataclass
from enum import Enum
import hashlib
import os
from pathlib import Path

from src.shared.kernel import Result, Success, Failure, Mono


class SupportedFileType(str, Enum):
    """지원되는 파일 형식"""

    PDF = "pdf"
    PPTX = "pptx"
    DOCX = "docx"
    TXT = "txt"
    MD = "md"
    MARKDOWN = "markdown"


@dataclass
class FileInfo:
    """파일 정보"""

    filename: str
    file_path: str
    file_size: int
    file_hash: str
    mime_type: str
    extension: str

    @classmethod
    def from_path(cls, file_path: str) -> "FileInfo":
        """파일 경로로부터 FileInfo 생성"""
        path = Path(file_path)

        # 파일 해시 계산
        file_hash = cls._calculate_file_hash(file_path)

        return cls(
            filename=path.name,
            file_path=str(path.absolute()),
            file_size=path.stat().st_size,
            file_hash=file_hash,
            mime_type=cls._get_mime_type(path.suffix),
            extension=path.suffix.lower(),
        )

    @staticmethod
    def _calculate_file_hash(file_path: str) -> str:
        """파일 해시 계산 (SHA-256)"""
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception:
            # 해시 계산 실패 시 파일명과 크기 기반 해시
            fallback_data = f"{file_path}_{os.path.getsize(file_path)}".encode()
            return hashlib.sha256(fallback_data).hexdigest()

    @staticmethod
    def _get_mime_type(extension: str) -> str:
        """확장자로부터 MIME 타입 추정"""
        mime_types = {
            ".pdf": "application/pdf",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".txt": "text/plain",
            ".md": "text/markdown",
            ".markdown": "text/markdown",
        }
        return mime_types.get(extension.lower(), "application/octet-stream")


@dataclass
class ExtractedContent:
    """추출된 텍스트 내용"""

    text: str
    metadata: Dict[str, Any]
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    char_count: Optional[int] = None

    def __post_init__(self) -> None:
        if self.word_count is None:
            self.word_count = len(self.text.split())
        if self.char_count is None:
            self.char_count = len(self.text)


class FileExtractor(ABC):
    """파일 추출기 기본 클래스"""

    @abstractmethod
    def extract_text(self, file_info: FileInfo) -> Mono[Result[ExtractedContent, str]]:
        """파일에서 텍스트 추출"""
        pass

    @abstractmethod
    def supports_file_type(self, file_type: SupportedFileType) -> bool:
        """파일 형식 지원 여부"""
        pass

    def validate_file(self, file_info: FileInfo) -> Result[None, str]:
        """파일 유효성 검증"""
        if not os.path.exists(file_info.file_path):
            return Failure(f"File not found: {file_info.file_path}")

        if file_info.file_size == 0:
            return Failure("Empty file")

        if file_info.file_size > 50 * 1024 * 1024:  # 50MB 제한
            return Failure("File too large (max 50MB)")

        return Success(None)


class TextFileExtractor(FileExtractor):
    """텍스트 파일 추출기 (TXT, MD)"""

    def supports_file_type(self, file_type: SupportedFileType) -> bool:
        return file_type in [
            SupportedFileType.TXT,
            SupportedFileType.MD,
            SupportedFileType.MARKDOWN,
        ]

    def extract_text(self, file_info: FileInfo) -> Mono[Result[ExtractedContent, str]]:
        """텍스트 파일에서 텍스트 추출"""

        async def extract() -> Result[ExtractedContent, str]:
            # 파일 검증
            validation_result = self.validate_file(file_info)
            if validation_result.is_failure():
                return Failure(validation_result.get_error() or "Validation failed")

            try:
                # 인코딩 감지 및 텍스트 읽기
                text = self._read_text_file(file_info.file_path)

                # 메타데이터 생성
                metadata = {
                    "extractor": "TextFileExtractor",
                    "encoding": "utf-8",  # 실제로는 감지된 인코딩
                    "file_type": file_info.extension,
                    "filename": file_info.filename,
                    "file_size": file_info.file_size,
                    "file_hash": file_info.file_hash,
                }

                content = ExtractedContent(text=text, metadata=metadata)

                return Success(content)

            except Exception as e:
                return Failure(f"Text extraction failed: {str(e)}")

        return Mono(lambda: asyncio.run(extract()))

    def _read_text_file(self, file_path: str) -> str:
        """텍스트 파일 읽기 (인코딩 자동 감지)"""
        encodings = ["utf-8", "utf-8-sig", "cp949", "euc-kr", "latin-1"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        # 모든 인코딩 실패 시 바이너리 읽기 후 에러 무시
        with open(file_path, "rb") as f:
            return f.read().decode("utf-8", errors="ignore")


class PDFExtractor(FileExtractor):
    """PDF 파일 추출기"""

    def supports_file_type(self, file_type: SupportedFileType) -> bool:
        return file_type == SupportedFileType.PDF

    def extract_text(self, file_info: FileInfo) -> Mono[Result[ExtractedContent, str]]:
        """PDF 파일에서 텍스트 추출"""

        async def extract() -> Result[ExtractedContent, str]:
            # 파일 검증
            validation_result = self.validate_file(file_info)
            if validation_result.is_failure():
                return Failure(validation_result.get_error() or "Validation failed")

            try:
                # PDF 텍스트 추출
                text, page_count = self._extract_pdf_text(file_info.file_path)

                # 메타데이터 생성
                metadata = {
                    "extractor": "PDFExtractor",
                    "library": "PyPDF2",  # 실제 사용된 라이브러리
                    "file_type": "pdf",
                    "filename": file_info.filename,
                    "file_size": file_info.file_size,
                    "file_hash": file_info.file_hash,
                }

                content = ExtractedContent(
                    text=text, metadata=metadata, page_count=page_count
                )

                return Success(content)

            except Exception as e:
                return Failure(f"PDF extraction failed: {str(e)}")

        return Mono(lambda: asyncio.run(extract()))

    def _extract_pdf_text(self, file_path: str) -> tuple[str, int]:
        """PDF에서 텍스트 추출 (PyPDF2 사용)"""
        try:
            # PyPDF2를 우선 시도
            import PyPDF2  # type: ignore

            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                # 파이프라인 패턴: PDF 텍스트 추출
                from src.shared.hof import pipe
                
                pdf_text_pipeline = pipe(
                    lambda reader: reader.pages,  # reader → pages
                    lambda pages: map(lambda page: page.extract_text(), pages),  # pages → text generator
                    list,  # generator → text_parts list
                    lambda text_parts: "\n\n".join(text_parts)  # text_parts → combined text
                )
                
                text = pdf_text_pipeline(reader)
                return text, len(reader.pages)

        except ImportError:
            try:
                # pymupdf를 fallback으로 시도
                import fitz  # type: ignore  # PyMuPDF

                doc = fitz.open(file_path)
                # 파이프라인 패턴: PyMuPDF 텍스트 추출
                pymupdf_text_pipeline = pipe(
                    lambda doc: map(lambda page: page.get_text(), doc),  # doc → text generator
                    list,  # generator → text_parts list
                    lambda text_parts: "\n\n".join(text_parts)  # text_parts → combined text
                )
                
                text = pymupdf_text_pipeline(doc)
                page_count = len(doc)
                doc.close()

                return text, page_count

            except ImportError:
                raise ImportError("PDF 추출을 위해 PyPDF2 또는 PyMuPDF가 필요합니다")


class PowerPointExtractor(FileExtractor):
    """PowerPoint 파일 추출기"""

    def supports_file_type(self, file_type: SupportedFileType) -> bool:
        return file_type == SupportedFileType.PPTX

    def extract_text(self, file_info: FileInfo) -> Mono[Result[ExtractedContent, str]]:
        """PPTX 파일에서 텍스트 추출"""

        async def extract() -> Result[ExtractedContent, str]:
            # 파일 검증
            validation_result = self.validate_file(file_info)
            if validation_result.is_failure():
                return Failure(validation_result.get_error() or "Validation failed")

            try:
                # PPTX 텍스트 추출
                text, slide_count = self._extract_pptx_text(file_info.file_path)

                # 메타데이터 생성
                metadata = {
                    "extractor": "PowerPointExtractor",
                    "library": "python-pptx",
                    "file_type": "pptx",
                    "filename": file_info.filename,
                    "file_size": file_info.file_size,
                    "file_hash": file_info.file_hash,
                }

                content = ExtractedContent(
                    text=text, metadata=metadata, page_count=slide_count
                )

                return Success(content)

            except Exception as e:
                return Failure(f"PowerPoint extraction failed: {str(e)}")

        return Mono(lambda: asyncio.run(extract()))

    def _extract_pptx_text(self, file_path: str) -> tuple[str, int]:
        """PPTX에서 텍스트 추출 (python-pptx 사용)"""
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(file_path)
            text_parts = []

            for i, slide in enumerate(prs.slides):
                slide_texts = []

                # 슬라이드의 모든 텍스트 추출 - HOF 패턴 사용
                from src.shared.hof import compact_map
                
                def extract_shape_text(shape):
                    if hasattr(shape, "text") and shape.text.strip():
                        return shape.text.strip()
                    return None
                
                slide_texts = compact_map(extract_shape_text, slide.shapes)

                if slide_texts:
                    slide_content = f"=== 슬라이드 {i+1} ===\n" + "\n".join(slide_texts)
                    text_parts.append(slide_content)

            text = "\n\n".join(text_parts)
            return text, len(prs.slides)

        except ImportError:
            raise ImportError("PPTX 추출을 위해 python-pptx가 필요합니다")


class WordExtractor(FileExtractor):
    """Word 문서 추출기"""

    def supports_file_type(self, file_type: SupportedFileType) -> bool:
        return file_type == SupportedFileType.DOCX

    def extract_text(self, file_info: FileInfo) -> Mono[Result[ExtractedContent, str]]:
        """DOCX 파일에서 텍스트 추출"""

        async def extract() -> Result[ExtractedContent, str]:
            # 파일 검증
            validation_result = self.validate_file(file_info)
            if validation_result.is_failure():
                return Failure(validation_result.get_error() or "Validation failed")

            try:
                # DOCX 텍스트 추출
                text, paragraph_count = self._extract_docx_text(file_info.file_path)

                # 메타데이터 생성
                metadata = {
                    "extractor": "WordExtractor",
                    "library": "python-docx",
                    "file_type": "docx",
                    "paragraph_count": paragraph_count,
                    "filename": file_info.filename,
                    "file_size": file_info.file_size,
                    "file_hash": file_info.file_hash,
                }

                content = ExtractedContent(text=text, metadata=metadata)

                return Success(content)

            except Exception as e:
                return Failure(f"Word extraction failed: {str(e)}")

        return Mono(lambda: asyncio.run(extract()))

    def _extract_docx_text(self, file_path: str) -> tuple[str, int]:
        """DOCX에서 텍스트 추출 (python-docx 사용)"""
        try:
            from docx import Document  # type: ignore

            doc = Document(file_path)
            # 파이프라인 패턴: Word 텍스트 추출
            from src.shared.hof import pipe
            
            docx_text_pipeline = pipe(
                lambda doc: doc.paragraphs,  # doc → paragraphs
                lambda paragraphs: filter(lambda p: p.text.strip(), paragraphs),  # → non-empty paragraphs
                lambda filtered: map(lambda p: p.text, filtered),  # → text generator
                list,  # generator → text_parts list
                lambda text_parts: "\n".join(text_parts)  # text_parts → combined text
            )
            
            text = docx_text_pipeline(doc)
            return text, len(doc.paragraphs)

        except ImportError:
            raise ImportError("DOCX 추출을 위해 python-docx가 필요합니다")
